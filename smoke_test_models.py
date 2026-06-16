"""
冒烟测试：校验百炼/OpenAI 兼容网关下各模型可调，以及一期默认主链路可用。

用法（在 xiangmushu 目录下）:
  python smoke_test_models.py
  python smoke_test_models.py --offline        # 不调外部 API：配置摘要 + 全链路离线断言（见下）
  python smoke_test_models.py --skip-vision   # 跳过多模态（省配额）
  python smoke_test_models.py --skip-chroma   # 跳过 Chroma+embedding
  python smoke_test_models.py --probe-models  # 网关 13 模型：chat/vision/enable_search 能力矩阵

--offline 默认门禁（与 docs/测试与验收.md 一期矩阵一致）:
  ContentGenerator 路由、query_expander、_max_output_tokens、
  WordFiller.clean_table_answer、锚点/占位相关 filler 离线回归。

说明：审核、task_grouper、evidence_planner、batch_generator、template_vision、
多模态表格等离线检查函数仍保留在脚本中供研发使用，但不再作为一期 `--offline` 默认门禁。

依赖：已配置 DASHSCOPE_API_KEY 或 OPENAI_API_KEY，且 .env 或环境变量可读。
"""
from __future__ import annotations

import argparse
import base64
import os
import re
import struct
import sys
import tempfile
import traceback
import zlib
from typing import Any, Dict, List, Optional, Tuple

# 复星网关可上架的对话模型（不含 embedding）
GATEWAY_PROBE_MODELS: List[str] = []

# 保证以仓库内模块方式加载
_ROOT = os.path.dirname(os.path.abspath(__file__))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

os.chdir(_ROOT)

# 先加载 .env
from dotenv import load_dotenv

load_dotenv(os.path.join(_ROOT, ".env"))

import config  # noqa: E402
from openai import OpenAI  # noqa: E402
from core.dashscope_chat import chat_completions_create, direct_chat_completions_create  # noqa: E402
from core.fill_task import FillTask  # noqa: E402
from core.generator import ContentGenerator, _max_output_tokens  # noqa: E402
from core.openai_embeddings import TimeoutOpenAIEmbedding  # noqa: E402


def _ordered_models(*models: str) -> List[str]:
    out: List[str] = []
    seen = set()
    for model in models:
        mid = (model or "").strip()
        if not mid or mid in seen:
            continue
        seen.add(mid)
        out.append(mid)
    return out


GATEWAY_PROBE_MODELS = _ordered_models(
    config.LARGE_LLM_MODEL,
    getattr(config, "FALLBACK_LLM_MODEL_1", ""),
    getattr(config, "FALLBACK_LLM_MODEL_2", ""),
    getattr(config, "FALLBACK_LLM_MODEL_3", ""),
    config.SMALL_LLM_MODEL,
    getattr(config, "SMALL_LLM_FALLBACK_MODEL", ""),
    config.VISION_WEB_MODEL,
    getattr(config, "VISION_WEB_FALLBACK_MODEL", ""),
    config.TEMPLATE_VISION_MODEL,
    getattr(config, "TEMPLATE_VISION_FALLBACK_MODEL", ""),
)


def _rgba_png_bytes(width: int, height: int) -> bytes:
    """生成纯色 RGBA PNG（百炼要求宽高均 > 10）。"""

    def _chunk(tag: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 6, 0, 0, 0)
    raw = b"".join(b"\x00" + b"\xff\x00\x00\xff" * width for _ in range(height))
    return (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", ihdr)
        + _chunk(b"IDAT", zlib.compress(raw, 9))
        + _chunk(b"IEND", b"")
    )


def _mask(s: str, keep: int = 6) -> str:
    s = (s or "").strip()
    if len(s) <= keep:
        return "(空或过短)"
    return s[:keep] + "…" + f"（共{len(s)}字符）"


def _client() -> OpenAI:
    if not config.chat_llm_configured():
        raise SystemExit("未配置聊天通道：请设置 DASHSCOPE_API_KEY 或 OPENAI_API_KEY。")
    return config.openai_client_for_chat()


def _print_config():
    print("=== 配置摘要 ===")
    print("  聊天主通道: 阿里云百炼 compatible-mode")
    print(f"  OPENAI_BASE_URL: {config.OPENAI_BASE_URL}")
    print(f"  EMBEDDING_OPENAI_BASE_URL: {config.EMBEDDING_OPENAI_BASE_URL}")
    print(f"  DASHSCOPE_COMPAT_BASE:     {config.DASHSCOPE_COMPAT_BASE}")
    print(f"  DASHSCOPE_API_KEY: {_mask(config.DASHSCOPE_API_KEY)}")
    print(f"  OPENAI_API_KEY:    {_mask(config.OPENAI_API_KEY)}")
    print(f"  嵌入/回落 Key:     {_mask(config.OPENAI_COMPAT_API_KEY)}")
    print(f"  SMALL_LLM_MODEL:   {config.SMALL_LLM_MODEL}  T={config.TEMP_SMALL_LLM}")
    print(
        f"  TEMPLATE_ANALYZE:  {config.TEMPLATE_ANALYZE_MODEL}  T={config.TEMP_TEMPLATE_ANALYZE}"
    )
    print(f"  LARGE_LLM_MODEL:   {config.LARGE_LLM_MODEL}  T={config.TEMP_LARGE_LLM}")
    print(f"  AUDIT_LLM_MODEL:   {config.AUDIT_LLM_MODEL}  T={config.TEMP_AUDIT}")
    print(f"  VISION_WEB_MODEL:  {config.VISION_WEB_MODEL}  T={config.TEMP_VISION}")
    print(
        f"  TEMPLATE_VISION:   {config.TEMPLATE_VISION_MODEL} · pages={config.TEMPLATE_VISION_MAX_PAGES} "
        f"zoom={config.TEMPLATE_VISION_ZOOM} long_edge={config.TEMPLATE_VISION_MAX_LONG_EDGE} "
        f"api_timeout={config.TEMPLATE_VISION_API_TIMEOUT}s enabled={config.TEMPLATE_VISION_ENABLED}"
    )
    print(f"  VISION_EXTRACT:    {config.VISION_EXTRACT_MODEL}")
    print(f"  TABLE_CELL_VISION: {config.TABLE_CELL_VISION_MODEL}")
    print(f"  EMBEDDING_MODEL:   {config.EMBEDDING_MODEL}")
    print()


def _chat_ping(client: OpenAI, model: str, temperature: float, label: str) -> bool:
    print(f"=== {label} ({model}) ===")
    try:
        r = chat_completions_create(
            client,
            model=model,
            messages=[{"role": "user", "content": "只回复一个字：好"}],
            temperature=temperature,
            max_tokens=16,
        )
        text = (r.choices[0].message.content or "").strip()
        print(f"  回复: {text!r}")
        if not text:
            print("  [FAIL] 空回复")
            return False
        print("  [OK]")
        return True
    except Exception as e:
        print(f"  [FAIL] {e}")
        if os.getenv("SMOKE_VERBOSE"):
            traceback.print_exc()
        return False


def _vision_data_url() -> str:
    png = _rgba_png_bytes(16, 16)
    return "data:image/png;base64," + base64.standard_b64encode(png).decode("ascii")


def _probe_one_gateway(
    client: OpenAI, model: str, kind: str
) -> Tuple[bool, str]:
    """仅走当前 client（不触发 dashscope_chat 回落），测网关原生能力。"""
    try:
        kwargs: dict[str, Any] = {
            "model": model,
            "temperature": 0.1 if kind == "chat" else 0.25,
            "max_tokens": 16 if kind == "chat" else (64 if kind == "vision" else 32),
        }
        if kind == "chat":
            kwargs["messages"] = [{"role": "user", "content": "只回复一个字：好"}]
        elif kind == "vision":
            kwargs["messages"] = [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "请用一句话描述这张图（若空白则说明）。"},
                        {"type": "image_url", "image_url": {"url": _vision_data_url()}},
                    ],
                }
            ]
        else:
            kwargs["messages"] = [
                {
                    "role": "user",
                    "content": "今天是星期几？只回答星期几，不要解释。",
                }
            ]
            kwargs["extra_body"] = {"enable_search": True}
        r = direct_chat_completions_create(client, **kwargs)
        text = (r.choices[0].message.content or "").strip()
        if not text:
            return False, "空回复"
        return True, text[:50]
    except Exception as e:
        return False, str(e)[:120]


def _probe_one(
    client: OpenAI, model: str, kind: str
) -> Tuple[bool, str]:
    """kind: chat | vision | search。经 dashscope_chat（含回落），返回 (ok, detail)。"""
    try:
        if kind == "chat":
            r = chat_completions_create(
                client,
                model=model,
                messages=[{"role": "user", "content": "只回复一个字：好"}],
                temperature=0.1,
                max_tokens=16,
            )
            text = (r.choices[0].message.content or "").strip()
            if not text:
                return False, "空回复"
            return True, text[:40]
        if kind == "vision":
            r = chat_completions_create(
                client,
                model=model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "请用一句话描述这张图（若空白则说明）。",
                            },
                            {
                                "type": "image_url",
                                "image_url": {"url": _vision_data_url()},
                            },
                        ],
                    }
                ],
                temperature=0.25,
                max_tokens=64,
            )
            text = (r.choices[0].message.content or "").strip()
            if not text:
                return False, "空回复"
            return True, text[:50]
        # search
        r = chat_completions_create(
            client,
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": "今天是星期几？只回答星期几，不要解释。",
                }
            ],
            temperature=0.2,
            max_tokens=32,
            extra_body={"enable_search": True},
        )
        text = (r.choices[0].message.content or "").strip()
        if not text:
            return False, "空回复"
        return True, text[:50]
    except Exception as e:
        return False, str(e)[:120]


def _run_probe_models(client: OpenAI, models: Optional[List[str]] = None) -> bool:
    """对网关模型跑 chat/vision/search 探测，打印 Markdown 矩阵。"""
    models = models or GATEWAY_PROBE_MODELS
    rows: List[Dict[str, Any]] = []
    print("=== 网关模型能力探测 ===")
    print(f"  模型数: {len(models)} · 通道: {_client_base_label()}")
    print()
    gw_only = False
    for model in models:
        row: Dict[str, Any] = {"model": model}
        for kind in ("chat", "vision", "search"):
            if gw_only:
                ok_gw, detail_gw = _probe_one_gateway(client, model, kind)
                ok_wrap, detail_wrap = _probe_one(client, model, kind)
                ok = ok_gw
                detail = detail_gw
                row[kind] = "OK" if ok_gw else "FAIL"
                row[f"{kind}_gw"] = row[kind]
                row[f"{kind}_wrap"] = "OK" if ok_wrap else "FAIL"
                if kind == "search" and (not ok_gw) and ok_wrap:
                    row[kind] = "FALLBACK"
                suffix = ""
                if kind == "search" and row.get("search_gw") == "FAIL" and row.get("search_wrap") == "OK":
                    suffix = " [仅回落百炼可用]"
                row[f"{kind}_detail"] = detail_gw
                mark = row[kind]
                print(f"  [{mark}] {model} | {kind}(网关): {detail_gw[:70]}{suffix}")
            else:
                ok, detail = _probe_one(client, model, kind)
                row[kind] = "OK" if ok else "FAIL"
                row[f"{kind}_detail"] = detail
                print(f"  [{'OK' if ok else 'FAIL'}] {model} | {kind}: {detail[:80]}")
        rows.append(row)
    print()
    if gw_only:
        print("| 模型 | chat(网关) | vision(网关) | search(网关) | search(含回落) |")
        print("|------|:----------:|:------------:|:------------:|:--------------:|")
        for row in rows:
            print(
                f"| {row['model']} | {row.get('chat_gw', row['chat'])} | "
                f"{row.get('vision_gw', row['vision'])} | {row.get('search_gw', row['search'])} | "
                f"{row.get('search_wrap', '-')} |"
            )
    else:
        print("| 模型 | chat | vision | enable_search |")
        print("|------|:----:|:------:|:-------------:|")
        for row in rows:
            print(
                f"| {row['model']} | {row['chat']} | {row['vision']} | {row['search']} |"
            )
    out_path = os.path.join(_ROOT, "data", "probe_gateway_models.md")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("# 阿里云百炼模型能力探测\n\n")
        f.write(f"通道: {_client_base_label()}\n\n")
        if gw_only:
            f.write(
                "说明：search(网关) 为直连网关，但同样统一关闭深度思考；"
                "search(含回落) 经 dashscope_chat（失败时用百炼 + VISION_WEB_MODEL）。\n\n"
            )
            f.write(
                "| 模型 | chat(网关) | vision(网关) | search(网关) | search(含回落) | 备注 |\n"
            )
            f.write("|------|:----------:|:------------:|:------------:|:--------------:|------|\n")
            for row in rows:
                note = ""
                if row.get("search_gw") == "FAIL" and row.get("search_wrap") == "OK":
                    note = "联网仅百炼回落"
                f.write(
                    f"| {row['model']} | {row.get('chat_gw')} | {row.get('vision_gw')} | "
                    f"{row.get('search_gw')} | {row.get('search_wrap')} | {note} |\n"
                )
        else:
            f.write("| 模型 | chat | vision | enable_search | chat 备注 | search 备注 |\n")
            f.write("|------|:----:|:------:|:-------------:|-----------|-------------|\n")
            for row in rows:
                f.write(
                    f"| {row['model']} | {row['chat']} | {row['vision']} | {row['search']} | "
                    f"{row.get('chat_detail', '')[:40]} | {row.get('search_detail', '')[:40]} |\n"
                )
    print(f"\n  矩阵已写入: {out_path}")
    chat_ok = sum(1 for r in rows if r["chat"] == "OK")
    return chat_ok >= 1


def _client_base_label() -> str:
    return f"阿里云百炼 {config.OPENAI_BASE_URL}"


def _vision_ping(client: OpenAI, model: Optional[str] = None) -> bool:
    model = model or config.VISION_WEB_MODEL
    print(f"=== 视觉模型 ({model}) ===")
    try:
        url = _vision_data_url()
        r = chat_completions_create(
            client,
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请用一句话描述这张图里有什么（若几乎空白则说明）。",
                        },
                        {"type": "image_url", "image_url": {"url": url}},
                    ],
                }
            ],
            temperature=config.TEMP_VISION,
            max_tokens=128,
        )
        text = (r.choices[0].message.content or "").strip()
        print(f"  回复片段: {text[:120]!r}…")
        if not text:
            print("  [FAIL] 空回复")
            return False
        print("  [OK]")
        return True
    except Exception as e:
        print(f"  [FAIL] {e}")
        if os.getenv("SMOKE_VERBOSE"):
            traceback.print_exc()
        return False


def _embedding_ping() -> bool:
    print(f"=== Embedding ({config.EMBEDDING_MODEL}) ===")
    try:
        fn = TimeoutOpenAIEmbedding(
            api_key=config.OPENAI_COMPAT_API_KEY or "sk-placeholder",
            base_url=config.EMBEDDING_OPENAI_BASE_URL or None,
            model_name=config.EMBEDDING_MODEL,
            timeout=config.OPENAI_TIMEOUT,
            max_retries=config.OPENAI_MAX_RETRIES,
        )
        vecs = fn(["测试向量入库用的短句。"])
        if not vecs or len(vecs[0]) < 8:
            print("  [FAIL] 向量维度过短或空")
            return False
        print(f"  向量维度: {len(vecs[0])}")
        print("  [OK]")
        return True
    except Exception as e:
        print(f"  [FAIL] {e}")
        if os.getenv("SMOKE_VERBOSE"):
            traceback.print_exc()
        return False


class _MockVS:
    """模拟向量库：控制 weak_kb 与检索命中。"""

    def __init__(self, count: int, results: List[dict[str, Any]]):
        self._count = count
        self._results = results

    def get_collection_count(self) -> int:
        return self._count

    def search(self, *args: Any, **kwargs: Any) -> List[dict[str, Any]]:
        return list(self._results)

    def get_all_documents(self, *args: Any, **kwargs: Any) -> List[dict[str, Any]]:
        return list(self._results)


def _routing_tests() -> bool:
    print("=== ContentGenerator 路由（无真实检索/联网）===")
    ok = True
    task = FillTask(
        task_id="t1",
        target_chapter="测试章",
        task_type="paragraph",
        description="写一句项目背景",
        location_hint={},
        word_limit=50,
    )

    # 有 KB 命中且相似度足够 -> 默认走小模型省 token（关 USE_SMALL_LLM_FOR_STRONG_RAG 则仍用大模型）
    g1 = ContentGenerator(_MockVS(5, [{"text": "参考资料", "metadata": {"source": "x"}, "distance": 0.3}]))
    _, m1, t1, _, rm1, _ = g1._build_chat_request(task, top_k=3, enable_web=True, retrieval_max_distance=1.0)
    exp1 = (
        config.SMALL_LLM_MODEL
        if config.USE_SMALL_LLM_FOR_STRONG_RAG
        else config.LARGE_LLM_MODEL
    )
    print(f"  有命中+高相似+联网开关(未触发弱库联网): model={m1!r} tier={rm1.get('generation_tier')}")
    if m1 != exp1:
        print(f"  [FAIL] 预期模型 {exp1!r}")
        ok = False
    else:
        print("  [OK] 路由模型符合预期")

    # 空库 + 开联网 -> VISION_WEB_MODEL + extra_body.enable_search（百炼内置联网）
    g2 = ContentGenerator(_MockVS(0, []))
    _, m2, t2, eb2, _, _ = g2._build_chat_request(task, top_k=3, enable_web=True, retrieval_max_distance=1.0)
    print(f"  空库 + 开联网: model={m2!r} temp={t2}")
    if m2 != config.VISION_WEB_MODEL:
        print("  [FAIL] 预期 VISION_WEB_MODEL（弱库联网档）")
        ok = False
    else:
        print("  [OK] 使用联网档（VISION_WEB_MODEL）")
    if not eb2.get("enable_search"):
        print("  [FAIL] 预期 extra_body.enable_search=True")
        ok = False
    else:
        print("  [OK] enable_search 已开启")

    # 空库 + 不开联网 -> 大模型
    g3 = ContentGenerator(_MockVS(0, []))
    _, m3, t3, eb3, _, _ = g3._build_chat_request(task, top_k=3, enable_web=False, retrieval_max_distance=1.0)
    print(f"  空库 + 关联网: model={m3!r} temp={t3}")
    if m3 != config.LARGE_LLM_MODEL:
        print("  [FAIL] 预期 LARGE_LLM_MODEL")
        ok = False
    else:
        print("  [OK] 使用大模型")
    if eb3.get("enable_search"):
        print("  [FAIL] 关联网不应设置 enable_search")
        ok = False

    # 有命中但 distance 大 -> 估算相似度低 + 开联网 -> 仍走联网档
    g4 = ContentGenerator(
        _MockVS(
            5,
            [{"text": "弱相关片段", "metadata": {"source": "x"}, "distance": 0.86}],
        )
    )
    _, m4, t4, eb4, rm4, _ = g4._build_chat_request(
        task, top_k=3, enable_web=True, retrieval_max_distance=1.0
    )
    print(f"  有命中但低相似 + 开联网: model={m4!r} low_sim={rm4.get('low_similarity')}")
    if config.FULL_RECALL_MODE:
        if rm4.get("full_recall_mode") is not True:
            print("  [FAIL] 全量召回开启时应标记 full_recall_mode=True")
            ok = False
        else:
            print("  [OK] 全量召回开启时跳过低相似联网回退")
    else:
        if m4 != config.VISION_WEB_MODEL:
            print("  [FAIL] 预期 VISION_WEB_MODEL（低相似度联网档）")
            ok = False
        else:
            print("  [OK] 低相似触发联网档")
        if not eb4.get("enable_search"):
            print("  [FAIL] 预期 extra_body.enable_search=True")
            ok = False

    # 联网档 + 创意模式：系统提示切换，用户提示不出现「请注明资料未载明」硬性句
    msgs_cr, _, _, _, rm_cr, _ = g2._build_chat_request(
        task,
        top_k=3,
        enable_web=True,
        retrieval_max_distance=1.0,
        web_writing_mode="creative",
    )
    if not rm_cr.get("web_creative_prompt"):
        print("  [FAIL] 弱库联网+创意应 web_creative_prompt=True")
        ok = False
    elif "联网创意" not in (msgs_cr[0].get("content") or ""):
        print("  [FAIL] 创意联网应使用联网创意系统提示")
        ok = False
    elif "无依据的要点请注明「资料未载明」" in (msgs_cr[1].get("content") or ""):
        print("  [FAIL] 创意用户提示不应含硬性资料未载明句")
        ok = False
    else:
        print("  [OK] 联网创意提示词")

    msgs_ca, _, _, _, rm_ca, _ = g2._build_chat_request(
        task,
        top_k=3,
        enable_web=True,
        retrieval_max_distance=1.0,
        web_writing_mode="calm",
    )
    if rm_ca.get("web_creative_prompt"):
        print("  [FAIL] 冷静模式 web_creative_prompt 应为 False")
        ok = False
    elif "无依据的要点请注明「资料未载明」" not in (msgs_ca[1].get("content") or ""):
        print("  [FAIL] 冷静用户提示应保留资料未载明要求")
        ok = False
    else:
        print("  [OK] 联网冷静提示词")

    return ok


def _audit_offline_helpers() -> bool:
    """不调 API：审核修订策略与字数上限。"""
    print("=== 审核辅助（offline）===")
    from core.content_auditor import AuditResult, effective_word_cap, should_apply_revision
    from core.fill_task import FillTask

    t_cell = FillTask(
        task_id="a1",
        target_chapter="表",
        task_type="table_cell",
        description="填一格",
        location_hint={"table_index": 0, "row": 1, "col": 1},
        word_limit=300,
    )
    if effective_word_cap(t_cell) != 120:
        print("  [FAIL] table_cell 字数上限应为 120")
        return False
    ar_long = AuditResult(
        verdict="minor_fix",
        issues=[],
        revised_content="x" * 300,
        one_line_summary="",
    )
    if should_apply_revision(t_cell, ar_long):
        print("  [FAIL] 过长修订稿不应采纳")
        return False
    ar_ok = AuditResult(
        verdict="minor_fix",
        issues=[],
        revised_content="符合资料",
        one_line_summary="",
    )
    if not should_apply_revision(t_cell, ar_ok):
        print("  [FAIL] 短修订稿应可采纳")
        return False
    print("  [OK]")
    return True


def _offline_query_expander() -> bool:
    print("=== query_expander（offline）===")
    from core.query_expander import expand_query

    q = expand_query("第一章", "项目实施背景与资金风险分析", "paragraph")
    if len(q) > 256:
        print(f"  [FAIL] 扩展 query 超长 len={len(q)}")
        return False
    if "项目" in q and "工程" not in q and "方案" not in q:
        print(f"  [FAIL] 预期含项目相关扩展词，得 {q!r}")
        return False
    q2 = expand_query("表", "技术路线简述", "table_cell")
    if len(q2) > 256:
        print("  [FAIL] table_cell query 超长")
        return False
    print("  [OK]")
    return True


def _offline_rule_and_need_model_audit() -> bool:
    print("=== rule_audit / need_model_audit（offline）===")
    from core.content_auditor import need_model_audit, rule_audit

    t_para = FillTask(
        task_id="p1",
        target_chapter="章",
        task_type="paragraph",
        description="简述",
        location_hint={},
        word_limit=300,
    )
    if len(rule_audit(t_para, "")) == 0:
        print("  [FAIL] 空答案应产生 rule_issues")
        return False
    issues = rule_audit(t_para, "以下是正文内容很长" + "x" * 10)
    if not any("禁用" in x or "前缀" in x for x in issues):
        print(f"  [FAIL] 禁用前缀未检出 issues={issues}")
        return False

    relaxed = rule_audit(
        t_para,
        "未提供详细数据说明" + "x" * 20,
        {"web_creative_prompt": True},
    )
    if any("未提供" in x for x in relaxed):
        print(f"  [FAIL] 创意联网下不应因「未提供」前缀拦截 relaxed={relaxed}")
        return False

    t_cell = FillTask(
        task_id="c1",
        target_chapter="表",
        task_type="table_cell",
        description="填格",
        location_hint={"table_index": 0, "row": 0, "col": 1},
        word_limit=40,
    )
    ni = rule_audit(t_cell, "第一行\n第二行")
    if not ni or "换行" not in ni[0]:
        print(f"  [FAIL] 表格换行未检出 {ni}")
        return False

    meta_web = {"native_web_search": True, "generation_tier": "large", "best_similarity_est": 0.9}
    if not need_model_audit(t_para, meta_web, []):
        print("  [FAIL] 联网应触发模型审核")
        return False
    meta_safe = {
        "native_web_search": False,
        "generation_tier": "small_rag",
        "best_similarity_est": 0.7,
    }
    if need_model_audit(t_para, meta_safe, []):
        print("  [FAIL] 低风险场景不应强制模型审核")
        return False
    t_policy = FillTask(
        task_id="p2",
        target_chapter="章",
        task_type="paragraph",
        description="政策合规要点",
        location_hint={},
        word_limit=100,
    )
    if not need_model_audit(t_policy, meta_safe, []):
        print("  [FAIL] 描述含高风险词应触发模型审核")
        return False
    print("  [OK]")
    return True


def _offline_max_output_tokens() -> bool:
    print("=== _max_output_tokens（offline）===")
    if _max_output_tokens(999, "table_cell") != 180:
        print("  [FAIL] table_cell 应为 180")
        return False
    v50 = _max_output_tokens(50, "paragraph")
    if not (256 <= v50 <= 1024):
        print(f"  [FAIL] 短段 50 字异常 {v50}")
        return False
    v300 = _max_output_tokens(300, "paragraph")
    if v300 > 1024:
        print(f"  [FAIL] 300 字段 cap 1024 违反 {v300}")
        return False
    v800 = _max_output_tokens(800, "paragraph")
    if v800 > int(config.GEN_MAX_TOKENS_HARD_CAP):
        print(f"  [FAIL] 长段超过硬顶 {v800}")
        return False
    if v800 < 512:
        print(f"  [FAIL] 长段 floor 512 违反 {v800}")
        return False
    print("  [OK]")
    return True


def _offline_task_grouper() -> bool:
    print("=== task_grouper（offline）===")
    from core.task_grouper import group_tasks

    tasks = [
        FillTask(
            task_id="a",
            target_chapter="同一章",
            task_type="table_cell",
            description="格1",
            location_hint={"table_index": 0, "row": 1, "col": 0},
            word_limit=20,
        ),
        FillTask(
            task_id="b",
            target_chapter="同一章",
            task_type="table_cell",
            description="格2",
            location_hint={"table_index": 0, "row": 1, "col": 1},
            word_limit=20,
        ),
        FillTask(
            task_id="c",
            target_chapter="同一章",
            task_type="paragraph",
            description="段1",
            location_hint={},
            word_limit=100,
        ),
    ]
    groups = group_tasks(tasks)
    if len(groups) != 2:
        print(f"  [FAIL] 预期 2 组（同行表+段落）得 {len(groups)}")
        return False
    row_group = next(g for g in groups if g.is_table_group)
    if len(row_group.tasks) != 2:
        print("  [FAIL] 同行应 2 个任务")
        return False
    print("  [OK]")
    return True


def _offline_evidence_planner() -> bool:
    print("=== evidence_planner（offline）===")
    from core.evidence_planner import compress_evidence, format_evidence, retrieve_for_group
    from core.task_grouper import TaskGroup

    grp = TaskGroup(
        group_id="table_0_row_1",
        tasks=[],
        shared_query="项目 风险 应对措施",
        table_index=0,
    )
    t = FillTask(
        task_id="t",
        target_chapter="风险",
        task_type="paragraph",
        description="项目风险与应对措施",
        location_hint={},
        word_limit=80,
    )
    ev = retrieve_for_group(
        _MockVS(
            3,
            [
                {
                    "text": "无关句子。项目风险主要包括技术风险。应对措施已明确。",
                    "distance": 0.25,
                },
                {"text": "另一段。", "distance": 0.4},
            ],
        ),
        grp,
        top_k=3,
        max_distance=1.0,
    )
    if ev.weak_kb or ev.kb_hits < 1:
        print("  [FAIL] 应有 KB 命中")
        return False
    comp = compress_evidence(ev, t, max_chars=200)
    if len(comp) > 200 or "风险" not in comp:
        print(f"  [FAIL] 压缩结果异常 {comp!r}")
        return False
    fmt = format_evidence(ev, max_chars=500)
    if len(fmt) < 5:
        print(f"  [FAIL] format_evidence 异常 {fmt!r}")
        return False
    print("  [OK]")
    return True


def _offline_bundle_evidence_route_parity() -> bool:
    """prepare_bundle_from_evidence 与 _build_chat_request 在相同检索语义下模型/联网档一致。"""
    print("=== prepare_bundle_from_evidence 路由对齐（offline）===")
    from core.evidence_planner import Evidence

    task = FillTask(
        task_id="t1",
        target_chapter="测试章",
        task_type="paragraph",
        description="写一句项目背景",
        location_hint={},
        word_limit=50,
    )

    def _cmp(vs: _MockVS, enable_web: bool, ev: Evidence, label: str) -> bool:
        g = ContentGenerator(vs)
        _, m1, _, eb1, rm1, _ = g._build_chat_request(
            task, top_k=3, enable_web=enable_web, retrieval_max_distance=1.0
        )
        b2 = g.prepare_bundle_from_evidence(task, ev, enable_web=enable_web)
        if m1 != b2.model:
            print(f"  [FAIL] {label} model 不一致 build={m1!r} bundle={b2.model!r}")
            return False
        if bool(eb1.get("enable_search")) != bool(b2.extra_body.get("enable_search")):
            print(f"  [FAIL] {label} enable_search 不一致")
            return False
        if rm1.get("generation_tier") != b2.route_meta.get("generation_tier"):
            print(
                f"  [FAIL] {label} tier 不一致 {rm1.get('generation_tier')} vs {b2.route_meta.get('generation_tier')}"
            )
            return False
        if rm1.get("use_small_llm_for_rag") != b2.route_meta.get("use_small_llm_for_rag"):
            print(f"  [FAIL] {label} use_small 不一致")
            return False
        return True

    hits = [{"text": "参考资料", "metadata": {"source": "x"}, "distance": 0.3}]
    vs1 = _MockVS(5, hits)
    ev1 = Evidence(
        group_id="g1",
        raw_results=list(hits),
        kb_hits=1,
        weak_kb=False,
        best_similarity=0.7,
    )
    if not _cmp(vs1, True, ev1, "强命中"):
        return False

    vs2 = _MockVS(0, [])
    ev2 = Evidence(group_id="g2", raw_results=[], kb_hits=0, weak_kb=True, best_similarity=None)
    if not _cmp(vs2, True, ev2, "弱库"):
        return False

    low_hits = [{"text": "弱相关", "metadata": {"source": "x"}, "distance": 0.86}]
    vs3 = _MockVS(5, low_hits)
    ev3 = Evidence(
        group_id="g3",
        raw_results=list(low_hits),
        kb_hits=1,
        weak_kb=False,
        best_similarity=0.14,
    )
    if not _cmp(vs3, True, ev3, "低相似"):
        return False

    t_cell = FillTask(
        task_id="tc",
        target_chapter="表",
        task_type="table_cell",
        description="填一格",
        location_hint={"table_index": 0, "row": 0, "col": 0},
        word_limit=40,
    )
    g4 = ContentGenerator(_MockVS(5, hits))
    ev4 = Evidence(
        group_id="g4",
        raw_results=list(hits),
        kb_hits=1,
        weak_kb=False,
        best_similarity=0.7,
    )
    b4 = g4.prepare_bundle_from_evidence(t_cell, ev4, enable_web=False, table_context="表头:名称")
    user = b4.messages[1]["content"]
    if isinstance(user, list):
        user = "\n".join(
            str(x.get("text", ""))
            for x in user
            if isinstance(x, dict) and x.get("type") == "text"
        )
    if "表格填写任务" not in user:
        print("  [FAIL] 表格任务 user 应含表格填写任务")
        return False
    print("  [OK]")
    return True


def _offline_filler_clean_table() -> bool:
    print("=== WordFiller.clean_table_answer（offline）===")
    from core.filler import WordFiller

    raw = "答案：\n\n某值"
    out = WordFiller.clean_table_answer(raw, word_limit=40)
    if out.startswith("答案"):
        print(f"  [FAIL] 前缀未剥离 {out!r}")
        return False
    long_in = "x" * 500
    short = WordFiller.clean_table_answer(long_in, word_limit=20)
    if len(short) > 90:
        print(f"  [FAIL] 超长未截断 len={len(short)}")
        return False
    print("  [OK]")
    return True


def _offline_batch_json_fence() -> bool:
    print("=== batch_generator._strip_json_fence（offline）===")
    import json

    from core import batch_generator as bg

    raw = '```json\n{"0": "甲", "1": "乙"}\n```'
    data = json.loads(bg._strip_json_fence(raw))
    if data.get("0") != "甲":
        print(f"  [FAIL] 解析失败 {data}")
        return False
    print("  [OK]")
    return True


def _offline_filler_paragraph_placeholder_only() -> bool:
    print("=== WordFiller.placeholder_only 段落（offline）===")
    import os
    import tempfile

    from docx import Document as Doc

    from core.fill_task import FillTask
    from core.filler import WordFiller

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    d.add_paragraph("第一章 测试章节")
    d.add_paragraph("说明文字请填写此处结尾")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    try:
        f = WordFiller()
        f.fill_template(
            path,
            [
                FillTask(
                    task_id="p1",
                    target_chapter="第一章 测试章节",
                    task_type="paragraph",
                    description="填空",
                    location_hint={
                        "replace_mode": "placeholder_only",
                        "paragraph_text": "说明文字",
                    },
                    word_limit=50,
                )
            ],
            ["答案"],
            out,
        )
        d2 = Doc(out)
        paras = [p.text for p in d2.paragraphs]
        body = paras[1] if len(paras) > 1 else ""
        if "说明文字" not in body or "结尾" not in body:
            print(f"  [FAIL] 说明或结尾丢失 {body!r}")
            return False
        if "答案" not in body or "请填写" in body:
            print(f"  [FAIL] 占位未替换或仍含「请填写」 {body!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_filler_abstract_hint_line() -> bool:
    print("=== WordFiller 摘要空段+提示行（offline）===")
    import os
    import tempfile

    from docx import Document as Doc
    from docx.enum.style import WD_STYLE_TYPE

    from core.fill_task import FillTask
    from core.filler import WordFiller

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    try:
        d.styles["Heading 1"]
    except KeyError:
        d.styles.add_style("Heading 1", WD_STYLE_TYPE.PARAGRAPH)
    d.add_paragraph("摘  要", style="Heading 1")
    d.add_paragraph("")
    d.add_paragraph("（请在此填写摘要正文）")
    d.add_paragraph("关键词：测试")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    try:
        WordFiller().fill_template(
            path,
            [
                FillTask(
                    task_id="abs",
                    target_chapter="摘  要",
                    task_type="paragraph",
                    description="摘要正文",
                    location_hint={"paragraph_text": "请在此填写"},
                    word_limit=500,
                )
            ],
            ["摘要：\n\n这是生成的摘要正文。"],
            out,
        )
        d2 = Doc(out)
        texts = [p.text for p in d2.paragraphs]
        joined = "\n".join(texts)
        if "请在此填写" in joined:
            print(f"  [FAIL] 仍含提示行 {texts!r}")
            return False
        if "这是生成的摘要正文" not in joined:
            print(f"  [FAIL] 未写入摘要 {texts!r}")
            return False
        body_para = ""
        for t in texts:
            if "这是生成的摘要正文" in t:
                body_para = t.strip()
                break
        if re.match(r"^摘\s*要", body_para):
            print(f"  [FAIL] 正文段不应再以「摘要」起头 {body_para!r}")
            return False
        if texts[2] if len(texts) > 2 else "" == "这是生成的摘要正文。":
            pass
        elif "这是生成的摘要正文" not in (texts[1] if len(texts) > 1 else ""):
            print(f"  [FAIL] 摘要应在提示行位置而非空段 {texts!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_filler_abstract_instruction_below_fill() -> bool:
    """「摘要：在以下填写…」指引行应被替换/清除，勿仅在下方空段叠稿。"""
    print("=== WordFiller 摘要在以下填写指引（offline）===")
    import os
    import tempfile

    from docx import Document as Doc
    from docx.enum.style import WD_STYLE_TYPE

    from core.fill_task import FillTask
    from core.filler import WordFiller

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    try:
        d.styles["Heading 1"]
    except KeyError:
        d.styles.add_style("Heading 1", WD_STYLE_TYPE.PARAGRAPH)
    d.add_paragraph("摘  要", style="Heading 1")
    d.add_paragraph("摘要：在以下填写正文，字数 300–500 字。")
    d.add_paragraph("")
    d.add_paragraph("关键词：测试")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    try:
        WordFiller().fill_template(
            path,
            [
                FillTask(
                    task_id="abs_inst",
                    target_chapter="摘要",
                    task_type="paragraph",
                    description="摘要正文",
                    location_hint={},
                    word_limit=500,
                )
            ],
            ["本项目面向申报场景，实现了基于 RAG 的计划书半自动撰写与 Word 回填。"],
            out,
        )
        d2 = Doc(out)
        texts = [p.text for p in d2.paragraphs]
        joined = "\n".join(texts)
        if "在以下填写" in joined:
            print(f"  [FAIL] 仍含指引行 {texts!r}")
            return False
        if "本项目面向申报场景" not in joined:
            print(f"  [FAIL] 未写入摘要 {texts!r}")
            return False
        body_parts = [
            t.strip()
            for t in texts
            if "本项目面向申报场景" in t
        ]
        if not body_parts:
            print(f"  [FAIL] 未找到正文段 {texts!r}")
            return False
        if len(body_parts) > 1:
            print(f"  [FAIL] 正文不应拆成多段叠写 {texts!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_filler_abstract_heading_alias() -> bool:
    """任务章节写「摘要」、模板标题为「摘  要」且占位为【】时仍能命中并写入。"""
    print("=== WordFiller 摘要章节别名匹配（offline）===")
    import os
    import tempfile

    from docx import Document as Doc
    from docx.enum.style import WD_STYLE_TYPE

    from core.fill_task import FillTask
    from core.filler import WordFiller

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    try:
        d.styles["Heading 1"]
    except KeyError:
        d.styles.add_style("Heading 1", WD_STYLE_TYPE.PARAGRAPH)
    d.add_paragraph("摘  要", style="Heading 1")
    d.add_paragraph("")
    d.add_paragraph("【请在此填写摘要正文】")
    d.add_paragraph("关键词：x")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    try:
        WordFiller().fill_template(
            path,
            [
                FillTask(
                    task_id="abs2",
                    target_chapter="摘要",
                    task_type="paragraph",
                    description="摘要",
                    location_hint={},
                    word_limit=500,
                )
            ],
            ["别名匹配后的摘要正文。"],
            out,
        )
        d2 = Doc(out)
        joined = "\n".join(p.text for p in d2.paragraphs)
        if "【请在此填写" in joined or "请在此填写" in joined:
            print(f"  [FAIL] 占位未替换 {joined!r}")
            return False
        if "别名匹配后的摘要正文" not in joined:
            print(f"  [FAIL] 未写入 {joined!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_filler_abstract_writing_rubric() -> bool:
    print("=== WordFiller 摘要撰写要求条识别（offline）===")
    import os
    import tempfile

    from docx import Document as Doc
    from docx.enum.style import WD_STYLE_TYPE

    from core.fill_task import FillTask
    from core.filler import WordFiller

    rubric = (
        "撰写要求\n"
        "• 用300—500字概述智能体应用项目：面向什么真实场景、解决什么问题、目标用户是谁。\n"
        "• 说明综合使用的核心能力：角色设定、工作流、数据库、知识库、插件调用、测试与演示方式。\n"
        "• 概括最终成果：可运行入口、主要功能、典型演示效果、创新点与不足。\n"
        "• 摘要中避免只写“我学会了什么”，应突出“项目做成了什么、如何做成、效果如何”。"
    )
    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    try:
        d.styles["Heading 1"]
    except KeyError:
        d.styles.add_style("Heading 1", WD_STYLE_TYPE.PARAGRAPH)
    d.add_paragraph("摘  要", style="Heading 1")
    d.add_paragraph("")
    d.add_paragraph(rubric)
    d.add_paragraph("")
    d.add_paragraph("关键词：测试")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    try:
        WordFiller().fill_template(
            path,
            [
                FillTask(
                    task_id="rub",
                    target_chapter="摘要",
                    task_type="paragraph",
                    description="摘要正文",
                    location_hint={},
                    word_limit=500,
                )
            ],
            ["摘要成稿已正确替换模板中的写作说明条目。"],
            out,
        )
        d2 = Doc(out)
        joined = "\n".join(p.text for p in d2.paragraphs)
        if "用300—500字" in joined or "我学会了什么" in joined:
            print(f"  [FAIL] rubric 应被替换 {joined[:400]!r}…")
            return False
        if "摘要成稿已正确替换模板中的写作说明条目" not in joined:
            print(f"  [FAIL] 未写入正文 {joined!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_filler_abstract_rubric_in_table() -> bool:
    """撰写要求位于表格单元格（灰框）时清空表内说明，正文写在下方段落。"""
    print("=== WordFiller 摘要表内撰写要求（offline）===")
    import os
    import tempfile

    from docx import Document as Doc
    from docx.enum.style import WD_STYLE_TYPE

    from core.fill_task import FillTask
    from core.filler import WordFiller

    rubric = (
        "撰写要求\n"
        "• 用300—500字概述智能体应用项目：面向什么真实场景、解决什么问题、目标用户是谁。\n"
        "• 说明综合使用的核心能力：角色设定、工作流、数据库、知识库、插件调用、测试与演示方式。\n"
        "• 概括最终成果：可运行入口、主要功能、典型演示效果、创新点与不足。\n"
        "• 摘要中避免只写“我学会了什么”，应突出“项目做成了什么、如何做成、效果如何”。"
    )
    body = "表内 rubric 已清除，摘要正文写在灰框下方段落。"
    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    try:
        d.styles["Heading 1"]
    except KeyError:
        d.styles.add_style("Heading 1", WD_STYLE_TYPE.PARAGRAPH)
    d.add_paragraph("摘  要", style="Heading 1")
    tbl = d.add_table(rows=1, cols=1)
    tbl.cell(0, 0).text = rubric
    d.add_paragraph("")
    d.add_paragraph("关键词：测试")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    try:
        WordFiller().fill_template(
            path,
            [
                FillTask(
                    task_id="rub_tbl",
                    target_chapter="摘要",
                    task_type="paragraph",
                    description="摘要正文",
                    location_hint={},
                    word_limit=500,
                )
            ],
            [body],
            out,
        )
        d2 = Doc(out)
        joined_para = "\n".join(p.text for p in d2.paragraphs)
        joined_tables = "\n".join(
            cell.text
            for t in d2.tables
            for row in t.rows
            for cell in row.cells
        )
        joined = joined_para + "\n" + joined_tables
        if "用300—500字" in joined or "我学会了什么" in joined or "撰写要求" in joined:
            print(f"  [FAIL] 表内 rubric 应清除 {joined[:500]!r}…")
            return False
        if len(d2.tables) > 0:
            print(f"  [FAIL] 摘要 rubric 单格表应删除，仍有 {len(d2.tables)} 个表")
            return False
        if body not in joined_para:
            print(f"  [FAIL] 正文应在段落而非表内 {joined_para!r}")
            return False
        if "关键词：测试" not in joined_para:
            print(f"  [FAIL] 关键词行应保留 {joined_para!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_innovation_table_scan_tasks() -> bool:
    print("=== 创新点三列表扫槽（offline）===")
    from core.slot_scanner import scan_placeholder_slots

    p = "data/templates/智能体应用开发实践.docx"
    slots = [
        s
        for s in scan_placeholder_slots(p)
        if s.location_hint.get("table_index") == 35
    ]
    keys = {(s.location_hint["row"], s.location_hint["col"]) for s in slots}
    if not {(1, 0), (1, 1), (1, 2)}.issubset(keys):
        print(f"  [FAIL] 第 1 数据行三列均应有任务，得 {keys}")
        return False
    if len(keys) < 5:
        print(f"  [FAIL] 创新点表任务过少 {keys}")
        return False
    if not any(s.word_limit <= 60 for s in slots if s.location_hint["col"] == 0):
        print("  [FAIL] 创新点列应有较短字数上限")
        return False
    print("  [OK]")
    return True


def _offline_filler_bracket_fill_slot() -> bool:
    """【请在此填写…】应写入占位行，而非 6.x 小节标题。"""
    print("=== WordFiller bracket 占位行（offline）===")
    import os
    import tempfile

    from docx import Document as Doc

    from core.fill_task import FillTask
    from core.filler import WordFiller

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    d.add_paragraph("第6章 总结与反思")
    d.add_paragraph("6.1 核心收获")
    d.add_paragraph("围绕角色设定与工作流进行总结。")
    d.add_paragraph("【请在此填写核心收获】")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    body = "本项目通过角色设定与工作流串联知识库，实现可演示的智能体应用。"
    try:
        WordFiller().fill_template(
            path,
            [
                FillTask(
                    task_id="b1",
                    target_chapter="第6章 总结与反思",
                    task_type="paragraph",
                    description="核心收获",
                    location_hint={
                        "paragraph_text": "【请在此填写核心收获】",
                        "replace_mode": "full",
                    },
                    word_limit=500,
                )
            ],
            [body],
            out,
        )
        d2 = Doc(out)
        texts = [p.text.strip() for p in d2.paragraphs if p.text.strip()]
        if "【请在此填写核心收获】" in "\n".join(texts):
            print(f"  [FAIL] 占位行应被替换 {texts!r}")
            return False
        if body not in texts:
            print(f"  [FAIL] 正文未写入 {texts!r}")
            return False
        if texts[1] != "6.1 核心收获":
            print(f"  [FAIL] 小节标题被改写 {texts!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_reconcile_bracket_task() -> bool:
    print("=== reconcile_fill_tasks 补 bracket 任务（offline）===")
    from core.fill_task import FillTask
    from core.task_reconcile import reconcile_fill_tasks

    analyzer: list = []
    scanner = [
        FillTask(
            task_id="s1",
            target_chapter="第6章 总结与反思",
            task_type="paragraph",
            description="填写：请在此填写未来展望",
            location_hint={
                "paragraph_text": "【请在此填写未来展望】",
                "replace_mode": "full",
            },
            word_limit=500,
        )
    ]
    merged = reconcile_fill_tasks(analyzer, scanner)
    if len(merged) != 1:
        print(f"  [FAIL] 应补 1 条任务，得 {len(merged)}")
        return False
    if merged[0].location_hint.get("paragraph_text") != "【请在此填写未来展望】":
        print(f"  [FAIL] hint 异常 {merged[0].location_hint!r}")
        return False
    print("  [OK]")
    return True


def _offline_batch_json_recover() -> bool:
    print("=== batch_generator._parse_batch_json 宽松解析（offline）===")
    from core.batch_generator import _parse_batch_json

    raw = '说明\n{"0": "甲", "1": "乙"}'
    data = _parse_batch_json(raw, 2)
    if not data or data.get(0) != "甲" or data.get(1) != "乙":
        print(f"  [FAIL] {data!r}")
        return False
    print("  [OK]")
    return True


def _offline_filler_guidance_beats_empty() -> bool:
    print("=== WordFiller 说明段优先于空段（offline）===")
    import os
    import tempfile

    from docx import Document as Doc
    from docx.enum.style import WD_STYLE_TYPE

    from core.fill_task import FillTask
    from core.filler import WordFiller

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    try:
        d.styles["Heading 1"]
    except KeyError:
        d.styles.add_style("Heading 1", WD_STYLE_TYPE.PARAGRAPH)
    try:
        d.styles["Heading 2"]
    except KeyError:
        d.styles.add_style("Heading 2", WD_STYLE_TYPE.PARAGRAPH)
    d.add_paragraph("第1章 单元测试章", style="Heading 1")
    d.add_paragraph("1.1 小节", style="Heading 2")
    d.add_paragraph(
        "说明项目背景、现实痛点和应用价值。建议从「真实场景—现有问题—价值」三个层次展开，避免只写概念介绍。"
    )
    d.add_paragraph("")
    d.save(path)
    out = path.replace(".docx", "_out.docx")
    try:
        WordFiller().fill_template(
            path,
            [
                FillTask(
                    task_id="g1",
                    target_chapter="第1章 单元测试章",
                    task_type="paragraph",
                    description="写背景",
                    location_hint={},
                    word_limit=200,
                )
            ],
            ["已替换的正文段落。"],
            out,
        )
        d2 = Doc(out)
        texts = [p.text for p in d2.paragraphs]
        if any("说明项目背景" in t for t in texts):
            print(f"  [FAIL] 模板说明段应被替换 {texts!r}")
            return False
        if not any("已替换的正文段落" in t for t in texts):
            print(f"  [FAIL] 未写入生成内容 {texts!r}")
            return False
    finally:
        for p in (path, out):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_docx_typography() -> bool:
    print("=== docx_typography 宋体小四（offline）===")
    import os
    import tempfile

    from docx import Document as Doc
    from docx.oxml.ns import qn

    from core.docx_typography import SZ_BODY, apply_document_typography

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    d.add_paragraph("正文段落")
    d.save(path)
    try:
        apply_document_typography(d)
        d.save(path)
        d2 = Doc(path)
        run = d2.paragraphs[0].runs[0]
        rpr = run._r.rPr
        if rpr is None:
            print("  [FAIL] 无 rPr")
            return False
        sz = rpr.find(qn("w:sz"))
        if sz is None or sz.get(qn("w:val")) != str(SZ_BODY):
            print(f"  [FAIL] sz 非小四 {sz.get(qn('w:val')) if sz is not None else None}")
            return False
        fonts = rpr.find(qn("w:rFonts"))
        if fonts is None:
            print("  [FAIL] 无 rFonts")
            return False
        ea = fonts.get(qn("w:eastAsia")) or ""
        if "宋" not in ea and fonts.get(qn("w:ascii")) != "SimSun":
            print(f"  [FAIL] 字体非宋体 {ea}")
            return False
    finally:
        try:
            os.remove(path)
        except OSError:
            pass
    print("  [OK]")
    return True


def _offline_template_vision_and_filler_cell() -> bool:
    print("=== template_vision / filler 单元格样式（offline）===")
    import os
    import tempfile

    from docx import Document as Doc

    from core.fill_task import FillTask
    from core.filler import WordFiller
    from core.generator import ContentGenerator, format_template_vision_block
    from core.template_vision import (
        apply_chapter_hints_to_tasks,
        parse_vision_profile_json,
        pick_chapter_style_hint,
    )

    raw = (
        '{"layout_notes":"多表","fill_strategy":"表格短答","style_observations":"小四",'
        '"chapter_hints":[{"chapter_anchor":"第三章 核心","hint":"须体现联网检索步骤"}]}'
    )
    prof = parse_vision_profile_json(raw)
    if prof.get("error"):
        print(f"  [FAIL] parse {prof}")
        return False
    h = pick_chapter_style_hint("第三章 核心能力实现过程", prof)
    if "联网" not in h and "检索" not in h:
        print(f"  [FAIL] chapter hint 未命中: {h!r}")
        return False
    tasks = [
        FillTask(
            task_id="t1",
            target_chapter="第三章 核心能力实现过程",
            task_type="paragraph",
            description="写一段",
            location_hint={},
            word_limit=100,
        )
    ]
    apply_chapter_hints_to_tasks(tasks, prof)
    if "chapter_style_hint" not in tasks[0].location_hint:
        print("  [FAIL] 未写入 chapter_style_hint")
        return False
    vb = format_template_vision_block(tasks[0])
    if "视觉摘要" not in vb or "表格短答" not in vb:
        print(f"  [FAIL] vision block {vb!r}")
        return False

    class _ZVS:
        def get_collection_count(self) -> int:
            return 0

        def search(self, *a, **k):
            return []

    g = ContentGenerator(_ZVS())
    _, _, _, _, _, _ = g._build_chat_request(
        tasks[0], top_k=2, enable_web=False, retrieval_max_distance=1.0
    )
    b = g.prepare_generation_bundle(tasks[0], top_k=2, enable_web=False)
    u = b.messages[1]["content"]
    if isinstance(u, list):
        u = "\n".join(
            str(x.get("text", ""))
            for x in u
            if isinstance(x, dict) and x.get("type") == "text"
        )
    if "视觉摘要" not in u:
        print("  [FAIL] 生成 user 未含视觉摘要")
        return False

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    d.add_paragraph("章节")
    tbl = d.add_table(rows=1, cols=1)
    tbl.cell(0, 0).text = "请填写"
    d.save(path)
    try:
        out = path.replace(".docx", "_out.docx")
        f = WordFiller()
        f.fill_template(
            path,
            [
                FillTask(
                    task_id="c1",
                    target_chapter="章节",
                    task_type="table_cell",
                    description="x",
                    location_hint={"table_index": 0, "row": 0, "col": 0},
                    word_limit=50,
                )
            ],
            ["已填"],
            out,
        )
        d2 = Doc(out)
        if "已填" not in d2.tables[0].cell(0, 0).text:
            print("  [FAIL] 单元格未写入")
            return False
    finally:
        for p in (path, path.replace(".docx", "_out.docx")):
            try:
                os.remove(p)
            except OSError:
                pass
    print("  [OK]")
    return True


def _offline_table_cell_multimodal_content() -> bool:
    print("=== build_table_cell_user_content 多模态（offline）===")
    from core.template_vision import build_table_cell_user_content

    tiny_png = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\xf8\x0f\x00"
        b"\x01\x05\x01\x02\xca\x89_\xeb\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    mm = build_table_cell_user_content("【表格填写任务】测试", [tiny_png])
    if not isinstance(mm, list) or len(mm) < 2:
        print(f"  [FAIL] 预期 list 多段 content，得 {type(mm)}")
        return False
    if mm[0].get("type") != "text" or mm[-1].get("type") != "image_url":
        print(f"  [FAIL] 段类型异常 {mm[0]} … {mm[-1]}")
        return False
    print("  [OK]")
    return True


def _offline_document_blocks_and_chunk_metadata() -> bool:
    print("=== parsed blocks / chunk metadata (offline) ===")
    import os
    import tempfile

    from docx import Document as Doc

    from core.chunker import Chunker
    from core.parser import DocumentParser
    from core.reporting import build_evidence_ref

    path = tempfile.mktemp(suffix=".docx")
    d = Doc()
    d.add_heading("第一章 项目概述", level=1)
    d.add_paragraph("这是项目概述正文。")
    tbl = d.add_table(rows=7, cols=2)
    tbl.cell(0, 0).text = "字段"
    tbl.cell(0, 1).text = "说明"
    long_cell = "这是一段较长的表格内容，用于验证长表格切块时会重复表头上下文。" * 10
    for i in range(1, 7):
        tbl.cell(i, 0).text = f"项{i}"
        tbl.cell(i, 1).text = long_cell
    d.save(path)

    try:
        parsed = DocumentParser().parse(path)
        if not parsed.blocks:
            print("  [FAIL] parse 后 blocks 为空")
            return False
        if parsed.blocks[0].block_type != "heading":
            print(f"  [FAIL] 首 block 类型异常: {parsed.blocks[0].block_type!r}")
            return False

        table_blocks = [block for block in parsed.blocks if block.block_type == "table"]
        if not table_blocks:
            print("  [FAIL] 未识别 table block")
            return False
        if "字段 | 说明" not in (table_blocks[0].table_header or ""):
            print(f"  [FAIL] table_header 异常: {table_blocks[0].table_header!r}")
            return False

        chunks = Chunker().chunk(parsed)
        if not chunks:
            print("  [FAIL] chunk 结果为空")
            return False
        seqs = [chunk.seq for chunk in chunks]
        if seqs != list(range(len(chunks))):
            print(f"  [FAIL] seq 不连续: {seqs!r}")
            return False

        first = chunks[0]
        required = {"source", "chapter", "type", "page", "seq", "source_type", "content_format", "ref_id"}
        if not required.issubset(first.metadata):
            print(f"  [FAIL] metadata 缺字段: {first.metadata!r}")
            return False
        ref = build_evidence_ref(first.metadata)
        if "seq=" not in ref or "page=" not in ref:
            print(f"  [FAIL] evidence ref 异常: {ref!r}")
            return False

        table_chunks = [chunk for chunk in chunks if chunk.metadata.get("type") == "table"]
        if len(table_chunks) < 2:
            print(f"  [FAIL] 预期长表格被拆成多块，实际只有 {len(table_chunks)} 块")
            return False
        if not all("字段 | 说明" in chunk.text for chunk in table_chunks[:2]):
            print("  [FAIL] 表格切块未保留表头上下文")
            return False
    finally:
        try:
            os.remove(path)
        except OSError:
            pass

    print("  [OK]")
    return True


def _offline_multiformat_parsed_blocks() -> bool:
    print("=== pdf / pptx / image parsed blocks (offline) ===")
    import os
    import tempfile
    from unittest.mock import patch

    from pypdf import PdfWriter
    from pptx import Presentation

    from core.kb_extract import path_to_parsed_document

    pdf_path = tempfile.mktemp(suffix=".pdf")
    pptx_path = tempfile.mktemp(suffix=".pptx")
    png_path = tempfile.mktemp(suffix=".png")

    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    with open(pdf_path, "wb") as f:
        writer.write(f)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "PPT 标题"
    slide.placeholders[1].text = "PPT 正文"
    prs.save(pptx_path)

    with open(png_path, "wb") as f:
        f.write(_rgba_png_bytes(16, 16))

    try:
        pdf_doc = path_to_parsed_document(pdf_path, original_name="sample.pdf")
        if pdf_doc.kb_source_type != "pdf" or not pdf_doc.blocks:
            print("  [FAIL] pdf 解析未返回 blocks")
            return False
        if pdf_doc.blocks[0].source_type != "pdf" or pdf_doc.blocks[0].page != 1:
            print(f"  [FAIL] pdf block metadata 异常: {pdf_doc.blocks[0]!r}")
            return False

        ppt_doc = path_to_parsed_document(pptx_path, original_name="deck.pptx")
        if ppt_doc.kb_source_type != "pptx" or not ppt_doc.blocks:
            print("  [FAIL] pptx 解析未返回 blocks")
            return False
        if "PPT 标题" not in ppt_doc.blocks[0].text or ppt_doc.blocks[0].source_type != "pptx":
            print(f"  [FAIL] pptx block 异常: {ppt_doc.blocks[0]!r}")
            return False

        with patch("core.vision_extract.describe_image_bytes", return_value="image summary"):
            image_doc = path_to_parsed_document(png_path, original_name="sample.png")
        if image_doc.kb_source_type != "image_vision" or not image_doc.blocks:
            print("  [FAIL] image 解析未返回 blocks")
            return False
        if image_doc.blocks[0].source_type != "image_vision" or "image summary" not in image_doc.blocks[0].text:
            print(f"  [FAIL] image block 异常: {image_doc.blocks[0]!r}")
            return False
    finally:
        for path in (pdf_path, pptx_path, png_path):
            try:
                os.remove(path)
            except OSError:
                pass

    print("  [OK]")
    return True


def _offline_quality_report_and_verifier() -> bool:
    print("=== quality report / verifier (offline) ===")
    import json
    import os
    import tempfile

    from docx import Document as Doc

    from core.post_fill_verifier import verify_filled_document
    from core.reporting import (
        build_generation_trace,
        build_quality_report,
        quality_report_summary,
        save_quality_report,
    )

    template_path = tempfile.mktemp(suffix=".docx")
    output_path = tempfile.mktemp(suffix=".docx")

    template_doc = Doc()
    template_doc.add_paragraph("项目申报书封面")
    template_table = template_doc.add_table(rows=1, cols=2)
    template_table.cell(0, 0).text = "评分项"
    template_table.cell(0, 1).text = "分值"
    template_doc.add_heading("第一章 项目概述", level=1)
    template_doc.save(template_path)

    output_doc = Doc()
    output_doc.add_paragraph("项目申报书封面")
    output_table = output_doc.add_table(rows=1, cols=2)
    output_table.cell(0, 0).text = "评分项"
    output_table.cell(0, 1).text = "分值"
    output_doc.add_heading("第一章 项目概述", level=1)
    output_doc.add_paragraph("这是已生成的正文。")
    output_doc.add_paragraph("{{待替换}}")
    output_doc.save(output_path)

    tasks = [
        FillTask(
            task_id="t-report",
            target_chapter="第二章 技术方案",
            task_type="paragraph",
            description="补充方案描述",
            location_hint={},
            word_limit=200,
        )
    ]

    try:
        checks = verify_filled_document(template_path, output_path, tasks)
        if checks["ok"]:
            print(f"  [FAIL] 预期校验失败，实际为 {checks!r}")
            return False
        if not checks["leftover_placeholders"]:
            print("  [FAIL] 未识别残留占位符")
            return False
        if "第二章 技术方案" not in checks["missing_chapters"]:
            print(f"  [FAIL] 未识别缺失章节: {checks['missing_chapters']!r}")
            return False

        trace = build_generation_trace(
            tasks[0],
            {
                "model": "demo-model",
                "generation_tier": "large",
                "kb_hits": 2,
                "evidence_refs": ["sample.docx | seq=3 | page=1 | docx"],
            },
            "这是已生成的正文。",
            audit_verdict="rule_issue",
            audit_issues=["内容过短"],
            revised=True,
        )
        report = build_quality_report(
            template_name="template.docx",
            output_path=output_path,
            traces=[trace],
            post_fill_checks=checks,
            visual_audit={"score": 81},
        )
        report_path = save_quality_report(output_path, report)
        summary = quality_report_summary(report)
        if "任务 1" not in summary or "残留占位 1" not in summary:
            print(f"  [FAIL] summary 异常: {summary!r}")
            return False

        with open(report_path, "r", encoding="utf-8") as f:
            payload = json.load(f)
        if payload.get("traces", [{}])[0].get("model") != "demo-model":
            print(f"  [FAIL] report 写入异常: {payload!r}")
            return False
        if payload.get("traces", [{}])[0].get("evidence_refs") != ["sample.docx | seq=3 | page=1 | docx"]:
            print(f"  [FAIL] evidence refs 写入异常: {payload!r}")
            return False
    finally:
        report_path = os.path.splitext(output_path)[0] + ".report.json"
        for path in (template_path, output_path, report_path):
            try:
                os.remove(path)
            except OSError:
                pass

    print("  [OK]")
    return True


def _run_all_offline() -> bool:
    steps = [
        ("路由", _routing_tests),
        ("query_expander", _offline_query_expander),
        ("max_output_tokens", _offline_max_output_tokens),
        ("clean_table_answer", _offline_filler_clean_table),
        ("filler_placeholder_only_para", _offline_filler_paragraph_placeholder_only),
        ("filler_abstract_hint", _offline_filler_abstract_hint_line),
        ("filler_abstract_instruction", _offline_filler_abstract_instruction_below_fill),
        ("filler_abstract_heading_alias", _offline_filler_abstract_heading_alias),
        ("filler_abstract_writing_rubric", _offline_filler_abstract_writing_rubric),
        ("filler_abstract_rubric_in_table", _offline_filler_abstract_rubric_in_table),
        ("innovation_table_scan", _offline_innovation_table_scan_tasks),
        ("filler_bracket_slot", _offline_filler_bracket_fill_slot),
        ("filler_guidance_beats_empty", _offline_filler_guidance_beats_empty),
        ("document_blocks_chunk_metadata", _offline_document_blocks_and_chunk_metadata),
        ("multiformat_parsed_blocks", _offline_multiformat_parsed_blocks),
        ("quality_report_verifier", _offline_quality_report_and_verifier),
    ]
    ok_all = True
    for name, fn in steps:
        try:
            if not fn():
                ok_all = False
        except Exception as e:
            print(f"  [FAIL] {name} 异常: {e}")
            if os.getenv("SMOKE_VERBOSE"):
                traceback.print_exc()
            ok_all = False
    return ok_all


def _stream_smoke(client: OpenAI) -> bool:
    """流式 + extra_body 路径（主生成同款）。"""
    print(f"=== 流式大模型 ({config.LARGE_LLM_MODEL}) ===")
    try:
        stream = chat_completions_create(
            client,
            model=config.LARGE_LLM_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": "请用一句话说明：什么是「项目申报书」（不超过40字）。",
                }
            ],
            temperature=config.TEMP_LARGE_LLM,
            max_tokens=32,
            stream=True,
        )
        buf: list[str] = []
        for chunk in stream:
            ch = chunk.choices[0] if chunk.choices else None
            if ch and ch.delta and ch.delta.content:
                buf.append(ch.delta.content)
        text = "".join(buf).strip()
        print(f"  拼接: {text!r}")
        if not text:
            print("  [FAIL] 流式无内容")
            return False
        print("  [OK]")
        return True
    except Exception as e:
        print(f"  [FAIL] {e}")
        if os.getenv("SMOKE_VERBOSE"):
            traceback.print_exc()
        return False


def _chroma_minimal_embed() -> bool:
    """独立临时目录：add + query.search（走 embed_query）。"""
    print("=== Chroma 持久化 + add + search（embedding / embed_query）===")
    try:
        from core.vector_store import VectorStore
        from core.chunker import Chunk
        import uuid

        d = tempfile.mkdtemp(prefix="smoke_chroma_")
        slug = f"smoke_{uuid.uuid4().hex[:8]}"
        vs = VectorStore(persist_dir=d, kb_slug=slug)
        ch = Chunk(
            id=f"id_{uuid.uuid4().hex[:12]}",
            text="冒烟测试片段，仅验证入库链路。",
            metadata={"source": "smoke_test_models.py"},
        )
        vs.add_documents([ch])
        n = vs.get_collection_count()
        print(f"  collection={vs.collection_name} count={n}")
        if n < 1:
            print("  [FAIL] count 未增长")
            return False
        hits = vs.search("冒烟测试", top_k=2, max_distance=2.5)
        print(f"  search 命中数: {len(hits)}")
        if not hits:
            print("  [FAIL] search 无结果（embed_query 或检索链异常）")
            return False
        print("  [OK]")
        return True
    except Exception as e:
        print(f"  [FAIL] {e}")
        if os.getenv("SMOKE_VERBOSE"):
            traceback.print_exc()
        return False


def main() -> int:
    ap = argparse.ArgumentParser(description="模型与链路冒烟测试")
    ap.add_argument(
        "--offline",
        action="store_true",
        help="不调外部 API：配置摘要 + 路由/审核/query/分组证据/bundle 对齐等离线断言",
    )
    ap.add_argument("--skip-vision", action="store_true", help="跳过多模态视觉请求")
    ap.add_argument("--skip-chroma", action="store_true", help="跳过 Chroma 入库")
    ap.add_argument("--skip-embed-api", action="store_true", help="跳过仅 embedding API（仍可做 chroma）")
    ap.add_argument(
        "--probe-models",
        action="store_true",
        help="探测 GATEWAY_PROBE_MODELS 的 chat/vision/enable_search 并输出矩阵",
    )
    ap.add_argument(
        "--leaderboard",
        action="store_true",
        help="运行项目 V/A/C/G 评分榜（见 scripts/run_leaderboard.py）",
    )
    ap.add_argument(
        "--leaderboard-quick",
        action="store_true",
        help="与 --leaderboard 联用：仅评测 quick_models 子集",
    )
    ap.add_argument(
        "--leaderboard-dry-run",
        action="store_true",
        help="与 --leaderboard 联用：不调 API",
    )
    args = ap.parse_args()

    _print_config()

    if args.leaderboard:
        from core.leaderboard.runner import run_leaderboard

        out = run_leaderboard(
            quick=args.leaderboard_quick,
            dry_run=args.leaderboard_dry_run,
            channels=["fosun", "dashscope"],
        )
        print(f"[OK] leaderboard 完成: {out}")
        return 0

    if args.probe_models:
        if not config.chat_llm_configured():
            print("[FAIL] 无可用聊天 API Key（网关或百炼）")
            return 1
        ok = _run_probe_models(_client())
        return 0 if ok else 1

    if args.offline:
        ok = _run_all_offline()
        print()
        print("=== 汇总（offline）===")
        print("  [OK] offline 检查通过" if ok else "  [FAIL] offline 检查失败")
        return 0 if ok else 1

    if not config.chat_llm_configured():
        print("[FAIL] 无可用聊天 API Key（网关或百炼）")
        return 1

    results: list[bool] = []
    client = _client()

    results.append(_chat_ping(client, config.SMALL_LLM_MODEL, config.TEMP_SMALL_LLM, "小模型"))
    results.append(_chat_ping(client, config.LARGE_LLM_MODEL, config.TEMP_LARGE_LLM, "大模型"))
    if not args.skip_vision:
        results.append(_vision_ping(client))
    else:
        print("=== 视觉模型 === [SKIP]")
        results.append(True)

    if not args.skip_embed_api:
        results.append(_embedding_ping())
    else:
        print("=== Embedding API === [SKIP]")
        results.append(True)

    results.append(_run_all_offline())
    results.append(_stream_smoke(client))

    if not args.skip_chroma:
        results.append(_chroma_minimal_embed())
    else:
        print("=== Chroma 入库 === [SKIP]")
        results.append(True)

    passed = sum(1 for x in results if x)
    total = len(results)
    print()
    print(f"=== 汇总: {passed}/{total} 通过 ===")
    return 0 if passed == total else 1


if __name__ == "__main__":
    sys.exit(main())

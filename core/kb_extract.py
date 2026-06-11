"""知识库多格式解析：docx / pdf / pptx / 图片（视觉模型）→ ParsedDocument。"""
from __future__ import annotations

import logging
import mimetypes
import os
from typing import Optional

from core.document_models import DocumentBlock
from core.parser import DocumentParser, ParsedDocument, Section


_LOG = logging.getLogger(__name__)

# 直接走专用解析器的扩展名
_DIRECT_EXTS = {
    ".docx", ".pdf", ".pptx",
    ".png", ".jpg", ".jpeg", ".webp", ".gif",
    ".md", ".markdown",
}

# MarkItDown 兜底的扩展名白名单
MARKITDOWN_FALLBACK_EXTS = {
    ".txt", ".csv", ".html", ".htm",
    ".xlsx", ".xls", ".doc",
}

SUPPORTED_KB_EXTS = _DIRECT_EXTS | MARKITDOWN_FALLBACK_EXTS


def is_supported_kb_extension(ext: str) -> bool:
    """返回该扩展名（小写）是否属于知识库支持的类型。"""
    return (ext or "").lower() in SUPPORTED_KB_EXTS


def _convert_with_markitdown(path: str) -> Optional[str]:
    """尝试用 MarkItDown 把文件转成 Markdown 文本。

    转换失败或返回空文本时返回 None 并记 warning，调用方据此走兜底。
    """
    try:
        from markitdown import MarkItDown
    except ImportError:
        _LOG.warning("markitdown not installed; cannot convert %s", path)
        return None
    try:
        md = MarkItDown()
        result = md.convert(path)
    except Exception as exc:
        _LOG.warning("markitdown.convert failed for %s: %s", path, exc)
        return None
    text = (getattr(result, "text_content", "") or "").strip()
    return text or None


def _synthetic_doc(
    filename: str,
    body: str,
    kb_source_type: str,
    *,
    blocks: Optional[list[DocumentBlock]] = None,
) -> ParsedDocument:
    text = (body or "").strip()
    if not text:
        text = "（未能提取到有效文本内容）"
    sec = Section(level=0, title="全文", content=text)
    use_blocks = list(blocks or [])
    if not use_blocks:
        use_blocks.append(
            DocumentBlock(
                text=text,
                page=1,
                block_type="text",
                source_type=kb_source_type,
                chapter=sec.title,
            )
        )
    return ParsedDocument(
        filename=filename,
        sections=[sec],
        raw_tables=[],
        kb_source_type=kb_source_type,
        blocks=use_blocks,
    )


def _extract_pdf_blocks_markitdown(path: str) -> list[DocumentBlock]:
    """Try MarkItDown for better PDF extraction (tables, formatting)."""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(path)
        text = result.text_content.strip()
        if not text:
            return []
        # Split by markdown headings or double newlines into blocks
        import re
        sections = re.split(r'\n(?=#{1,3}\s)', text)
        blocks: list[DocumentBlock] = []
        for idx, section in enumerate(sections, start=1):
            section = section.strip()
            if not section:
                continue
            # Extract heading if present
            heading_match = re.match(r'^(#{1,3})\s+(.+?)\n', section)
            chapter = heading_match.group(2) if heading_match else f"第{idx}节"
            blocks.append(
                DocumentBlock(
                    text=section,
                    page=idx,
                    block_type="text",
                    source_type="pdf_markitdown",
                    chapter=chapter,
                    metadata={"parser": "markitdown"},
                )
            )
        return blocks
    except Exception:
        return []


def _extract_pdf_blocks_pypdf(path: str) -> list[DocumentBlock]:
    """Fallback PDF extraction using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    blocks: list[DocumentBlock] = []
    for idx, page in enumerate(reader.pages, start=1):
        t = page.extract_text()
        if t:
            blocks.append(
                DocumentBlock(
                    text=t.strip(),
                    page=idx,
                    block_type="text",
                    source_type="pdf",
                    chapter=f"第{idx}页",
                    metadata={"page_label": idx},
                )
            )
    return blocks


def _extract_pdf_blocks(path: str) -> list[DocumentBlock]:
    """PDF extraction: try MarkItDown first, fallback to pypdf."""
    blocks = _extract_pdf_blocks_markitdown(path)
    if blocks:
        return blocks
    return _extract_pdf_blocks_pypdf(path)


def _extract_pptx_blocks(path: str) -> list[DocumentBlock]:
    from pptx import Presentation

    prs = Presentation(path)
    blocks: list[DocumentBlock] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text.strip())
        if lines:
            blocks.append(
                DocumentBlock(
                    text="\n".join(lines),
                    page=i,
                    block_type="slide",
                    source_type="pptx",
                    chapter=f"第{i}页",
                    metadata={"slide_index": i},
                )
            )
    return blocks


def _guess_mime(path: str) -> str:
    mime, _ = mimetypes.guess_type(path)
    if mime and mime.startswith("image/"):
        return mime
    ext = os.path.splitext(path)[1].lower()
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
        ".gif": "image/gif",
    }.get(ext, "application/octet-stream")


def _extract_markdown_blocks_from_text(text: str) -> list[DocumentBlock]:
    """把一段已读好的 Markdown 文本切分为 DocumentBlock 列表。"""
    import re

    if not text:
        return []
    sections = re.split(r"\n(?=#{1,6}\s)", text)
    blocks: list[DocumentBlock] = []
    for idx, section in enumerate(sections, start=1):
        chunk = section.strip()
        if not chunk:
            continue
        heading_match = re.match(r"^(#{1,6})\s+(.+?)\n", chunk)
        chapter = heading_match.group(2).strip() if heading_match else f"第{idx}节"
        blocks.append(
            DocumentBlock(
                text=chunk,
                page=idx,
                block_type="text",
                source_type="markdown",
                chapter=chapter,
                content_format="markdown",
                metadata={"parser": "markdown"},
            )
        )
    return blocks


def _extract_markdown_blocks(path: str) -> list[DocumentBlock]:
    with open(path, "r", encoding="utf-8") as f:
        text = f.read().strip()
    return _extract_markdown_blocks_from_text(text)


def path_to_parsed_document(path: str, original_name: Optional[str] = None) -> ParsedDocument:
    """
    根据扩展名解析为 ParsedDocument，供 Chunker 使用。
    original_name 用于展示与 metadata.source（默认取 path 的 basename）。
    """
    name = original_name or os.path.basename(path)
    ext = os.path.splitext(path)[1].lower()

    if ext == ".docx":
        return DocumentParser().parse(path)

    if ext == ".pdf":
        blocks = _extract_pdf_blocks(path)
        text = "\n\n".join(block.text for block in blocks)
        if not (text or "").strip():
            text = (
                "（本 PDF 未提取到文本层，可能为纯扫描件。请导出为图片后上传，"
                "或使用带文字层的 PDF。）"
            )
        return _synthetic_doc(name, text, "pdf", blocks=blocks)

    if ext == ".pptx":
        blocks = _extract_pptx_blocks(path)
        text = "\n\n".join(
            f"【第{block.page}页】\n{block.text}" for block in blocks
        )
        if not (text or "").strip():
            text = "（未从 PPTX 提取到文本，可能幻灯片主要为图片；可导出为图片后上传以启用视觉解析。）"
        return _synthetic_doc(name, text, "pptx", blocks=blocks)

    if ext in (".png", ".jpg", ".jpeg", ".webp", ".gif"):
        from core.vision_extract import describe_image_bytes

        with open(path, "rb") as f:
            raw = f.read()
        mime = _guess_mime(path)
        text = describe_image_bytes(raw, mime)
        blocks = [
            DocumentBlock(
                text=text,
                page=1,
                block_type="image",
                source_type="image_vision",
                chapter="图片视觉解析",
                metadata={"mime_type": mime},
            )
        ]
        doc = _synthetic_doc(
            name,
            f"【图片视觉解析】{name}\n\n{text}",
            "image_vision",
            blocks=blocks,
        )
        return doc

    if ext in (".md", ".markdown"):
        blocks = _extract_markdown_blocks(path)
        text = "\n\n".join(block.text for block in blocks)
        if not (text or "").strip():
            text = "（Markdown 文件未提取到有效文本内容）"
        return _synthetic_doc(name, text, "markdown", blocks=blocks)

    if ext in MARKITDOWN_FALLBACK_EXTS:
        text = _convert_with_markitdown(path)
        if text:
            blocks = _extract_markdown_blocks_from_text(text)
            body = "\n\n".join(b.text for b in blocks) if blocks else text
            if not (body or "").strip():
                body = "（未从该文件提取到有效文本内容）"
            return _synthetic_doc(name, body, "markdown", blocks=blocks or None)

    raise ValueError(f"不支持的文件类型: {ext}")
